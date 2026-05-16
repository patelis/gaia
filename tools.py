import os
import json
import base64
from pathlib import Path

from dotenv import load_dotenv
load_dotenv()

from langchain_community.tools import DuckDuckGoSearchRun
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain_community.document_loaders import WikipediaLoader
from langchain_community.document_loaders import ArxivLoader
from langchain_core.tools import tool

from huggingface_hub import InferenceClient

from utils import load_config, load_prompt

_config = load_config()
_vlm_model_name = _config["models"]["vlm"]["model_name"]
_vlm_system_prompt = load_prompt("prompts/vlm_prompt.yaml").content
_asr_model_name = _config["models"]["asr"]["model_name"]
_hf_client = InferenceClient(token=os.getenv("HF_INFERENCE_KEY"))

_ddg_search = None
_tavily_search = None

def _get_ddg():
    global _ddg_search
    if _ddg_search is None:
        _ddg_search = DuckDuckGoSearchRun()
    return _ddg_search

def _get_tavily():
    global _tavily_search
    if _tavily_search is None:
        _tavily_search = TavilySearchResults(max_results=3)
    return _tavily_search

# ============================================
# Basic Tools
# ============================================

@tool
def calculator(a: float, b: float, type: str) -> float:
    """Performs mathematical calculations, addition, subtraction, multiplication, division, modulus.
    Args: 
        a (float): first float number
        b (float): second float number
        type (str): the type of calculation to perform, can be addition, subtraction, multiplication, division, modulus
    """

    if type == "addition":
        return a + b
    elif type == "subtraction":
        return a - b
    elif type == "multiplication":
        return a * b
    elif type == "division":
        if b == 0:
            raise ValueError("Cannot divide by zero.")
        return a / b
    elif type == "modulus":
        return a % b
    else:
        raise TypeError(f"{type} is not an option for type, choose one of addition, subtraction, multiplication, division, modulus")

@tool
def duck_web_search(query: str) -> str:
    """Use DuckDuckGo to search the web.

    Args:
        query: The search query.
    """
    search = _get_ddg().invoke(query=query)
    
    return {"duckduckgo_web_search": search}

@tool
def wiki_search(query: str) -> str:
    """Search Wikipedia for a query and return maximum 3 results.
    
    Args:
        query: The search query."""
    documents = WikipediaLoader(query=query, load_max_docs=3).load()
    processed_documents = "\n\n---\n\n".join(
        [
            f'Document title: {document.metadata.get("title", "")}. Summary: {document.metadata.get("summary", "")}. Documents details: {document.page_content}'
            for document in documents
        ])
    return {"wiki_results": processed_documents}

@tool
def arxiv_search(query: str) -> str:
    """Search Arxiv for a query and return maximum 3 result.
    
    Args:
        query: The search query."""
    documents = ArxivLoader(query=query, load_max_docs=3).load()
    processed_documents = "\n\n---\n\n".join(
        [
            f'Document title: {document.metadata.get("title", "")}. Summary: {document.metadata.get("summary", "")}. Documents details: {document.page_content}'
            for document in documents
        ])
    return {"arxiv_results": processed_documents}

@tool
def tavily_web_search(query: str) -> str:
    """Search the web using Tavily for a query and return maximum 3 results.
    
    Args:
        query: The search query."""
    search_documents = _get_tavily().invoke(input=query)
    web_results = "\n\n---\n\n".join(
        [
            f'Document title: {document["title"]}. Contents: {document["content"]}. Relevance Score: {document["score"]}'
            for document in search_documents
        ])
    return {"web_results": web_results}


@tool
def fetch_webpage(url: str) -> str:
    """
    Fetch and extract the main text content from a webpage.
    Use this when a search result points to a specific URL you need to read in full.

    Args:
        url: The full URL of the page to fetch.

    Returns:
        The extracted text content of the page.
    """
    import trafilatura
    try:
        downloaded = trafilatura.fetch_url(url)
        if downloaded is None:
            return f"Error: Could not fetch {url}"
        text = trafilatura.extract(downloaded, include_tables=True, include_links=False)
        if text is None:
            return f"Error: Could not extract content from {url}"
        return f"Page content from {url}:\n\n{text}"
    except Exception as e:
        return f"Error fetching webpage: {e}"


@tool
def python_eval(code: str) -> str:
    """
    Execute a Python code snippet and return its stdout output.
    Use this when a question asks what a script outputs, or when computation requires running code.

    Args:
        code: Python source code to execute.

    Returns:
        The stdout output of the code, or an error/timeout message.
    """
    import subprocess
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
            f.write(code)
            tmp_path = f.name
        result = subprocess.run(
            ['python3', tmp_path],
            capture_output=True, text=True, timeout=30
        )
        os.unlink(tmp_path)
        if result.returncode == 0:
            return f"Output:\n{result.stdout}"
        return f"Error (exit {result.returncode}):\n{result.stderr}"
    except subprocess.TimeoutExpired:
        return "Error: execution timed out (30s limit)"
    except Exception as e:
        return f"Error: {e}"


# ============================================
# VLM Tool
# ============================================

@tool
def analyze_image(image_path: str, question: str) -> str:
    """
    Analyze an image using a Vision Language Model (VLM) to answer a specific question.
    
    Args:
        image_path: Path to the image file (JPG, PNG).
        question: The specific question to answer about the image.
    
    Returns:
        A detailed description or answer based on the visual content.
    """
    try:
        if not os.path.exists(image_path):
            return f"Error: Image file not found at {image_path}"

        with open(image_path, "rb") as img_file:
            image_data = base64.b64encode(img_file.read()).decode("utf-8")
        ext = Path(image_path).suffix.lower().lstrip(".")
        mime_type = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
        image_url = f"data:{mime_type};base64,{image_data}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": image_url}},
                    {"type": "text", "text": f"{_vlm_system_prompt}\n\nQuestion: {question}"}
                ]
            }
        ]

        output = _hf_client.chat_completion(
            messages=messages,
            model=_vlm_model_name,
            max_tokens=1000
        )

        return output.choices[0].message.content

    except Exception as e:
        return f"Error analyzing image with VLM: {str(e)}"


# ============================================
# Document Processing Tools
# ============================================

@tool
def read_pdf(file_path: str) -> str:
    """
    Extract text content from a PDF file.
    
    Args:
        file_path: Path to the PDF file to read.
        
    Returns:
        The text content of the PDF, with page separators.
    """
    from pypdf import PdfReader
    
    try:
        reader = PdfReader(file_path)
        text = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text.append(f"--- Page {i+1} ---\n{page_text}")
        
        return "\n\n".join(text) if text else "[Empty PDF]"
    except Exception as e:
        return f"Error reading PDF: {e}"


@tool
def read_docx(file_path: str) -> str:
    """
    Extract text content from a Word document (.docx).
    
    Args:
        file_path: Path to the Word document to read.
        
    Returns:
        The text content of the document.
    """
    from docx import Document
    
    try:
        doc = Document(file_path)
        text_parts = []

        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        if paragraphs:
            text_parts.append("\n".join(paragraphs))

        for i, table in enumerate(doc.tables):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            rows = [r for r in rows if r.strip()]
            if rows:
                text_parts.append(f"--- Table {i+1} ---\n" + "\n".join(rows))

        return "\n\n".join(text_parts) if text_parts else "[Empty document]"
    except Exception as e:
        return f"Error reading DOCX: {e}"


@tool
def read_pptx(file_path: str) -> str:
    """
    Extract text content from a PowerPoint presentation (.pptx).
    
    Args:
        file_path: Path to the PowerPoint file to read.
        
    Returns:
        The text content from all slides.
    """
    from pptx import Presentation
    
    try:
        prs = Presentation(file_path)
        text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:
                text.append("\n".join(slide_text))
        
        return "\n\n".join(text) if text else "[Empty presentation]"
    except Exception as e:
        return f"Error reading PPTX: {e}"


@tool
def read_text_file(file_path: str) -> str:
    """
    Read content from a plain text file (.txt).
    
    Args:
        file_path: Path to the text file to read.
        
    Returns:
        The content of the text file.
    """
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        return f"Error reading text file: {e}"


# ============================================
# Data Processing Tools (using polars)
# ============================================

@tool
def read_csv(file_path: str) -> str:
    """
    Read and analyze a CSV file using polars.
    
    Args:
        file_path: Path to the CSV file to read.
        
    Returns:
        Summary of the CSV including schema, row count, and data preview.
    """
    import polars as pl
    
    try:
        df = pl.read_csv(file_path)
        
        output = f"CSV File — {len(df)} rows, {len(df.columns)} columns\n"
        output += f"Columns: {df.columns}\n\n"
        output += f"Column Statistics:\n{df.describe()}\n\n"
        output += f"Data (first 20 rows):\n{df.head(20)}"
        if len(df) <= 50:
            output += f"\n\nComplete data:\n{df}"
        return output
    except Exception as e:
        return f"Error reading CSV: {e}"


@tool
def read_excel(file_path: str, sheet_id: int = 0) -> str:
    """
    Read and analyze an Excel file (.xlsx) using polars.
    
    Args:
        file_path: Path to the Excel file to read.
        sheet_id: The sheet index to read (0-based). Default is 0 (first sheet).
        
    Returns:
        Summary of the Excel sheet including schema, row count, and data preview.
    """
    import polars as pl
    import openpyxl

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        sheet_names = wb.sheetnames
        wb.close()
    except Exception:
        sheet_names = []

    try:
        df = pl.read_excel(file_path, sheet_id=sheet_id)
        sheet_label = sheet_names[sheet_id] if sheet_id < len(sheet_names) else str(sheet_id)

        output = f"Excel File — Available sheets: {sheet_names}\n\n"
        output += f"Sheet {sheet_id} ('{sheet_label}') — {len(df)} rows, {len(df.columns)} columns\n"
        output += f"Columns: {df.columns}\n\n"
        output += f"Column Statistics:\n{df.describe()}\n\n"
        output += f"Data (first 20 rows):\n{df.head(20)}"
        if len(df) <= 50:
            output += f"\n\nComplete data:\n{df}"
        return output
    except Exception as e:
        return f"Error reading Excel: {e}"


@tool
def read_jsonld(file_path: str) -> str:
    """
    Read and parse a JSON-LD file.
    
    Args:
        file_path: Path to the JSON-LD file to read.
        
    Returns:
        The formatted JSON content.
    """
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
        return f"JSON-LD Content:\n{json.dumps(data, indent=2)}"
    except Exception as e:
        return f"Error reading JSON-LD: {e}"


@tool
def read_pdb(file_path: str) -> str:
    """
    Read and analyze a PDB (Protein Data Bank) file for protein structure analysis.
    
    Args:
        file_path: Path to the PDB file to read.
        
    Returns:
        Analysis of the protein structure including atoms, chains, and coordinates.
    """
    from Bio.PDB import PDBParser
    import numpy as np
    
    try:
        parser = PDBParser(QUIET=True)
        structure = parser.get_structure("protein", file_path)
        
        info = ["=== PDB Structure Analysis ==="]
        
        atoms = list(structure.get_atoms())
        info.append(f"Total atoms: {len(atoms)}")
        
        for model in structure:
            info.append(f"\nModel {model.id}:")
            for chain in model:
                residues = list(chain.get_residues())
                info.append(f"  Chain {chain.id}: {len(residues)} residues")
        
        if len(atoms) >= 2:
            info.append("\nFirst atoms (for distance calculations):")
            for i, atom in enumerate(atoms[:5]):
                coord = atom.get_coord()
                info.append(
                    f"  Atom {i+1}: {atom.get_name()} at "
                    f"[{coord[0]:.4f}, {coord[1]:.4f}, {coord[2]:.4f}]"
                )
            
            dist = np.linalg.norm(atoms[0].get_coord() - atoms[1].get_coord())
            info.append(f"\nDistance between first two atoms: {dist:.4f} Angstroms")
        
        return "\n".join(info)
    except Exception as e:
        return f"Error reading PDB: {e}"


# ============================================
# Audio Processing Tools
# ============================================

@tool
def transcribe_audio(file_path: str) -> str:
    """
    Transcribe an audio file (MP3, WAV, etc.) to text using Whisper.
    
    Args:
        file_path: Path to the audio file to transcribe.
        
    Returns:
        The transcribed text from the audio.
    """
    try:
        result = _hf_client.automatic_speech_recognition(audio=file_path, model=_asr_model_name)
        return f"Audio Transcription:\n{result.text}"
    except Exception as e:
        return f"Error transcribing audio: {e}"


# ============================================
# Code Processing Tools
# ============================================

@tool
def read_python_file(file_path: str) -> str:
    """
    Read a Python source code file.
    
    Args:
        file_path: Path to the Python file to read.
        
    Returns:
        The Python code content.
    """
    try:
        with open(file_path, 'r') as f:
            code = f.read()
        return f"Python Code:\n```python\n{code}\n```"
    except Exception as e:
        return f"Error reading Python file: {e}"


# ============================================
# Archive Processing Tools
# ============================================

@tool
def extract_zip(file_path: str) -> str:
    """
    Extract a ZIP archive and list its contents.
    
    Args:
        file_path: Path to the ZIP file to extract.
        
    Returns:
        List of files extracted from the archive with their paths.
    """
    import zipfile
    
    try:
        extract_dir = Path(file_path).parent / Path(file_path).stem
        extract_dir.mkdir(exist_ok=True)
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)
        
        results = [f"ZIP Archive extracted to: {extract_dir}\n\nContents:"]
        
        for root, dirs, files in os.walk(extract_dir):
            for file in files:
                full_path = os.path.join(root, file)
                rel_path = os.path.relpath(full_path, extract_dir)
                file_size = os.path.getsize(full_path)
                results.append(f"  - {rel_path} ({file_size} bytes)")
        
        results.append(f"\nUse the appropriate read tool on the extracted files at: {extract_dir}/")
        
        return "\n".join(results)
    except Exception as e:
        return f"Error extracting ZIP: {e}"


# ============================================
# Generic File Processing
# ============================================

@tool
def read_file(file_path: str) -> str:
    """
    Automatically read a file based on its extension.
    
    Supported formats: PDF, DOCX, PPTX, TXT, CSV, XLSX, JSON-LD, PDB, Python, ZIP, JPG, JPEG, PNG, MP3, WAV, FLAC, OGG, M4A
    
    Args:
        file_path: Path to the file to read.
        
    Returns:
        The processed content of the file.
    """
    ext = Path(file_path).suffix.lower()
    
    processors = {
        '.pdf': lambda p: read_pdf.invoke(p),
        '.docx': lambda p: read_docx.invoke(p),
        '.pptx': lambda p: read_pptx.invoke(p),
        '.txt': lambda p: read_text_file.invoke(p),
        '.csv': lambda p: read_csv.invoke(p),
        '.xlsx': lambda p: read_excel.invoke(p),
        '.jsonld': lambda p: read_jsonld.invoke(p),
        '.pdb': lambda p: read_pdb.invoke(p),
        '.py': lambda p: read_python_file.invoke(p),
        '.mp3': lambda p: transcribe_audio.invoke(p),
        '.wav': lambda p: transcribe_audio.invoke(p),
        '.flac': lambda p: transcribe_audio.invoke(p),
        '.ogg': lambda p: transcribe_audio.invoke(p),
        '.m4a': lambda p: transcribe_audio.invoke(p),
        '.zip': lambda p: extract_zip.invoke(p),
        '.jpg': lambda p: analyze_image.invoke({"image_path": p, "question": "Describe this image in detail."}),
        '.jpeg': lambda p: analyze_image.invoke({"image_path": p, "question": "Describe this image in detail."}),
        '.png': lambda p: analyze_image.invoke({"image_path": p, "question": "Describe this image in detail."}),
    }
    
    processor = processors.get(ext)
    
    if processor:
        return processor(file_path)
    
    return f"[Unsupported file type: {ext}]"

# ============================================
# List of all tools
# ============================================

tools_list = [
    calculator,
    duck_web_search,
    wiki_search,
    arxiv_search,
    tavily_web_search,
    fetch_webpage,
    python_eval,
    read_pdf,
    read_docx,
    read_pptx,
    read_text_file,
    read_csv,
    read_excel,
    read_jsonld,
    read_pdb,
    transcribe_audio,
    read_python_file,
    extract_zip,
    analyze_image,
    read_file,
]