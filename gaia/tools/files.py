"""File / document / archive reading tools."""
import json
import os
import zipfile
from pathlib import Path

from langchain_core.tools import tool


@tool
def read_pdf(file_path: str) -> str:
    """
    Extract text content from a PDF file.

    Args:
        file_path: Path to the PDF file to read.

    Returns:
        The text content of the PDF, with page separators.
    """
    from pypdf import PdfReader

    try:
        reader = PdfReader(file_path)
        text = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text.append(f"--- Page {i+1} ---\n{page_text}")

        return "\n\n".join(text) if text else "[Empty PDF]"
    except Exception as e:
        return f"[read_pdf] failed to read PDF: {e}"


@tool
def read_docx(file_path: str) -> str:
    """
    Extract text content from a Word document (.docx).

    Args:
        file_path: Path to the Word document to read.

    Returns:
        The text content of the document.
    """
    from docx import Document

    try:
        doc = Document(file_path)
        text_parts = []

        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        if paragraphs:
            text_parts.append("\n".join(paragraphs))

        for i, table in enumerate(doc.tables):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            rows = [r for r in rows if r.strip()]
            if rows:
                text_parts.append(f"--- Table {i+1} ---\n" + "\n".join(rows))

        return "\n\n".join(text_parts) if text_parts else "[Empty document]"
    except Exception as e:
        return f"[read_docx] failed to read DOCX: {e}"


@tool
def read_pptx(file_path: str) -> str:
    """
    Extract text content from a PowerPoint presentation (.pptx).

    Args:
        file_path: Path to the PowerPoint file to read.

    Returns:
        The text content from all slides.
    """
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:
                text.append("\n".join(slide_text))

        return "\n\n".join(text) if text else "[Empty presentation]"
    except Exception as e:
        return f"[read_pptx] failed to read PPTX: {e}"


@tool
def read_text_file(file_path: str) -> str:
    """
    Read content from a plain text file (.txt).

    Args:
        file_path: Path to the text file to read.

    Returns:
        The content of the text file.
    """
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        return f"[read_text_file] failed: {e}"


@tool
def read_csv(file_path: str) -> str:
    """
    Read and analyze a CSV file using polars.

    Args:
        file_path: Path to the CSV file to read.

    Returns:
        Summary of the CSV including schema, row count, and data preview.
    """
    import polars as pl

    try:
        df = pl.read_csv(file_path)

        output = f"CSV File — {len(df)} rows, {len(df.columns)} columns\n"
        output += f"Columns: {df.columns}\n\n"
        output += f"Column Statistics:\n{df.describe()}\n\n"
        output += f"Data (first 20 rows):\n{df.head(20)}"
        if len(df) <= 50:
            output += f"\n\nComplete data:\n{df}"
        return output
    except Exception as e:
        return f"[read_csv] failed to read CSV: {e}"


@tool
def read_excel(file_path: str, sheet_id: int = 0) -> str:
    """
    Read and analyze an Excel file (.xlsx) using polars.

    Args:
        file_path: Path to the Excel file to read.
        sheet_id: The sheet index to read (0-based). Default is 0 (first sheet).

    Returns:
        Summary of the Excel sheet including schema, row count, and data preview.
    """
    import polars as pl
    import openpyxl

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        sheet_names = wb.sheetnames
        wb.close()
    except Exception:
        sheet_names = []

    try:
        df = pl.read_excel(file_path, sheet_id=sheet_id)
        sheet_label = sheet_names[sheet_id] if sheet_id < len(sheet_names) else str(sheet_id)

        output = f"Excel File — Available sheets: {sheet_names}\n\n"
        output += f"Sheet {sheet_id} ('{sheet_label}') — {len(df)} rows, {len(df.columns)} columns\n"
        output += f"Columns: {df.columns}\n\n"
        output += f"Column Statistics:\n{df.describe()}\n\n"
        output += f"Data (first 20 rows):\n{df.head(20)}"
        if len(df) <= 50:
            output += f"\n\nComplete data:\n{df}"
        return output
    except Exception as e:
        return f"[read_excel] failed to read Excel: {e}"


@tool
def read_jsonld(file_path: str) -> str:
    """
    Read and parse a JSON-LD file.

    Args:
        file_path: Path to the JSON-LD file to read.

    Returns:
        The formatted JSON content.
    """
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
        return f"JSON-LD Content:\n{json.dumps(data, indent=2)}"
    except Exception as e:
        return f"[read_jsonld] failed to read JSON-LD: {e}"


@tool
def read_pdb(file_path: str) -> str:
    """
    Read and analyze a PDB (Protein Data Bank) file for protein structure analysis.

    Args:
        file_path: Path to the PDB file to read.

    Returns:
        Analysis of the protein structure including atoms, chains, and coordinates.
    """
    from Bio.PDB import PDBParser
    import numpy as np

    try:
        parser = PDBParser(QUIET=True)
        structure = parser.get_structure("protein", file_path)

        info = ["=== PDB Structure Analysis ==="]

        atoms = list(structure.get_atoms())
        info.append(f"Total atoms: {len(atoms)}")

        for model in structure:
            info.append(f"\nModel {model.id}:")
            for chain in model:
                residues = list(chain.get_residues())
                info.append(f"  Chain {chain.id}: {len(residues)} residues")

        if len(atoms) >= 2:
            info.append("\nFirst atoms (for distance calculations):")
            for i, atom in enumerate(atoms[:5]):
                coord = atom.get_coord()
                info.append(
                    f"  Atom {i+1}: {atom.get_name()} at "
                    f"[{coord[0]:.4f}, {coord[1]:.4f}, {coord[2]:.4f}]"
                )

            dist = np.linalg.norm(atoms[0].get_coord() - atoms[1].get_coord())
            info.append(f"\nDistance between first two atoms: {dist:.4f} Angstroms")

        return "\n".join(info)
    except Exception as e:
        return f"[read_pdb] failed to read PDB: {e}"


@tool
def read_python_file(file_path: str) -> str:
    """
    Read a Python source code file.

    Args:
        file_path: Path to the Python file to read.

    Returns:
        The Python code content.
    """
    try:
        with open(file_path, 'r') as f:
            code = f.read()
        return f"Python Code:\n```python\n{code}\n```"
    except Exception as e:
        return f"[read_python_file] failed: {e}"


@tool
def extract_zip(file_path: str) -> str:
    """
    Extract a ZIP archive and list its contents.

    Args:
        file_path: Path to the ZIP file to extract.

    Returns:
        List of files extracted from the archive with their paths.
    """
    try:
        extract_dir = Path(file_path).parent / Path(file_path).stem
        extract_dir.mkdir(exist_ok=True)

        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)

        results = [f"ZIP Archive extracted to: {extract_dir}\n\nContents:"]

        for root, dirs, files in os.walk(extract_dir):
            for file in files:
                full_path = os.path.join(root, file)
                rel_path = os.path.relpath(full_path, extract_dir)
                file_size = os.path.getsize(full_path)
                results.append(f"  - {rel_path} ({file_size} bytes)")

        results.append(f"\nUse the appropriate read tool on the extracted files at: {extract_dir}/")

        return "\n".join(results)
    except Exception as e:
        return f"[extract_zip] failed: {e}"
