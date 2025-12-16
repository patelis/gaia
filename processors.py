"""
File processing tools for the GAIA agent.

These tools allow the agent to process different file types on demand.
All processors are decorated with @tool for LangChain integration.

Supported extensions:
- Documents: .pdf, .docx, .pptx, .txt
- Data: .csv, .xlsx, .jsonld, .pdb
- Media: .jpg, .png, .mp3
- Code: .py
- Archives: .zip
"""

import json
import os
from pathlib import Path
from langchain_core.tools import tool


# ============================================
# Document Processing Tools
# ============================================

@tool
def read_pdf(file_path: str) -> str:
    """
    Extract text content from a PDF file.
    
    Args:
        file_path: Path to the PDF file to read.
        
    Returns:
        The text content of the PDF, with page separators.
    """
    from pypdf import PdfReader
    
    reader = PdfReader(file_path)
    text = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text:
            text.append(f"--- Page {i+1} ---\n{page_text}")
    
    return "\n\n".join(text) if text else "[Empty PDF]"


@tool
def read_docx(file_path: str) -> str:
    """
    Extract text content from a Word document (.docx).
    
    Args:
        file_path: Path to the Word document to read.
        
    Returns:
        The text content of the document.
    """
    from docx import Document
    
    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    
    return "\n".join(paragraphs) if paragraphs else "[Empty document]"


@tool
def read_pptx(file_path: str) -> str:
    """
    Extract text content from a PowerPoint presentation (.pptx).
    
    Args:
        file_path: Path to the PowerPoint file to read.
        
    Returns:
        The text content from all slides.
    """
    from pptx import Presentation
    
    prs = Presentation(file_path)
    text = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if len(slide_text) > 1:
            text.append("\n".join(slide_text))
    
    return "\n\n".join(text) if text else "[Empty presentation]"


@tool
def read_text_file(file_path: str) -> str:
    """
    Read content from a plain text file (.txt).
    
    Args:
        file_path: Path to the text file to read.
        
    Returns:
        The content of the text file.
    """
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


# ============================================
# Data Processing Tools (using polars)
# ============================================

@tool
def read_csv(file_path: str) -> str:
    """
    Read and analyze a CSV file using polars.
    
    Args:
        file_path: Path to the CSV file to read.
        
    Returns:
        Summary of the CSV including schema, row count, and data preview.
    """
    import polars as pl
    
    df = pl.read_csv(file_path)
    
    output = f"""CSV File Summary:
- Total rows: {len(df)}
- Columns: {df.columns}
- Schema: {dict(df.schema)}

Data (first 20 rows):
{df.head(20)}
"""
    
    if len(df) <= 50:
        output += f"\nComplete data:\n{df}"
    
    return output


@tool
def read_excel(file_path: str, sheet_id: int = 0) -> str:
    """
    Read and analyze an Excel file (.xlsx) using polars.
    
    Args:
        file_path: Path to the Excel file to read.
        sheet_id: The sheet index to read (0-based). Default is 0 (first sheet).
        
    Returns:
        Summary of the Excel sheet including schema, row count, and data preview.
    """
    import polars as pl
    
    df = pl.read_excel(file_path, sheet_id=sheet_id)
    
    output = f"""Excel Sheet {sheet_id} Summary:
- Total rows: {len(df)}
- Columns: {df.columns}

Data (first 20 rows):
{df.head(20)}
"""
    
    if len(df) <= 50:
        output += f"\nComplete data:\n{df}"
    
    return output


@tool
def read_jsonld(file_path: str) -> str:
    """
    Read and parse a JSON-LD file.
    
    Args:
        file_path: Path to the JSON-LD file to read.
        
    Returns:
        The formatted JSON content.
    """
    with open(file_path, 'r') as f:
        data = json.load(f)
    
    return f"JSON-LD Content:\n{json.dumps(data, indent=2)}"


@tool
def read_pdb(file_path: str) -> str:
    """
    Read and analyze a PDB (Protein Data Bank) file for protein structure analysis.
    
    Args:
        file_path: Path to the PDB file to read.
        
    Returns:
        Analysis of the protein structure including atoms, chains, and coordinates.
    """
    from Bio.PDB import PDBParser
    import numpy as np
    
    parser = PDBParser(QUIET=True)
    structure = parser.get_structure("protein", file_path)
    
    info = ["=== PDB Structure Analysis ==="]
    
    atoms = list(structure.get_atoms())
    info.append(f"Total atoms: {len(atoms)}")
    
    for model in structure:
        info.append(f"\nModel {model.id}:")
        for chain in model:
            residues = list(chain.get_residues())
            info.append(f"  Chain {chain.id}: {len(residues)} residues")
    
    if len(atoms) >= 2:
        info.append("\nFirst atoms (for distance calculations):")
        for i, atom in enumerate(atoms[:5]):
            coord = atom.get_coord()
            info.append(
                f"  Atom {i+1}: {atom.get_name()} at "
                f"[{coord[0]:.4f}, {coord[1]:.4f}, {coord[2]:.4f}]"
            )
        
        dist = np.linalg.norm(atoms[0].get_coord() - atoms[1].get_coord())
        info.append(f"\nDistance between first two atoms: {dist:.4f} Angstroms")
    
    return "\n".join(info)


# ============================================
# Audio Processing Tools
# ============================================

@tool
def transcribe_audio(file_path: str) -> str:
    """
    Transcribe an audio file (MP3, WAV, etc.) to text using Whisper.
    
    Args:
        file_path: Path to the audio file to transcribe.
        
    Returns:
        The transcribed text from the audio.
    """
    from transformers import pipeline
    
    transcriber = pipeline(
        "automatic-speech-recognition", 
        model="openai/whisper-base",
        chunk_length_s=30
    )
    
    result = transcriber(file_path)
    
    return f"Audio Transcription:\n{result['text']}"


# ============================================
# Code Processing Tools
# ============================================

@tool
def read_python_file(file_path: str) -> str:
    """
    Read a Python source code file.
    
    Args:
        file_path: Path to the Python file to read.
        
    Returns:
        The Python code content.
    """
    with open(file_path, 'r') as f:
        code = f.read()
    
    return f"Python Code:\n```python\n{code}\n```"


# ============================================
# Archive Processing Tools
# ============================================

@tool
def extract_zip(file_path: str) -> str:
    """
    Extract a ZIP archive and list its contents.
    
    Args:
        file_path: Path to the ZIP file to extract.
        
    Returns:
        List of files extracted from the archive with their paths.
    """
    import zipfile
    
    extract_dir = Path(file_path).parent / Path(file_path).stem
    extract_dir.mkdir(exist_ok=True)
    
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    
    results = [f"ZIP Archive extracted to: {extract_dir}\n\nContents:"]
    
    for root, dirs, files in os.walk(extract_dir):
        for file in files:
            full_path = os.path.join(root, file)
            rel_path = os.path.relpath(full_path, extract_dir)
            file_size = os.path.getsize(full_path)
            results.append(f"  - {rel_path} ({file_size} bytes)")
    
    results.append(f"\nUse the appropriate read tool on the extracted files at: {extract_dir}/")
    
    return "\n".join(results)


# ============================================
# Image Processing Tools
# ============================================

@tool
def describe_image(file_path: str) -> str:
    """
    Get basic information about an image file (JPG, PNG).
    
    Note: This provides metadata only. For visual content analysis,
    a vision model would be needed.
    
    Args:
        file_path: Path to the image file.
        
    Returns:
        Image metadata including format, size, and mode.
    """
    from PIL import Image
    
    img = Image.open(file_path)
    
    return f"""Image Information:
- File: {Path(file_path).name}
- Format: {img.format}
- Size: {img.size[0]} x {img.size[1]} pixels
- Mode: {img.mode}

[Note: For visual content analysis, examine the image directly or use a vision model.]
"""


# ============================================
# Generic File Processing
# ============================================

@tool
def read_file(file_path: str) -> str:
    """
    Automatically read a file based on its extension.
    
    Supported formats: PDF, DOCX, PPTX, TXT, CSV, XLSX, JSONLD, PDB, PY, MP3, ZIP, JPG, PNG
    
    Args:
        file_path: Path to the file to read.
        
    Returns:
        The processed content of the file.
    """
    ext = Path(file_path).suffix.lower()
    
    processors = {
        '.pdf': lambda p: read_pdf.invoke(p),
        '.docx': lambda p: read_docx.invoke(p),
        '.pptx': lambda p: read_pptx.invoke(p),
        '.txt': lambda p: read_text_file.invoke(p),
        '.csv': lambda p: read_csv.invoke(p),
        '.xlsx': lambda p: read_excel.invoke(p),
        '.jsonld': lambda p: read_jsonld.invoke(p),
        '.pdb': lambda p: read_pdb.invoke(p),
        '.py': lambda p: read_python_file.invoke(p),
        '.mp3': lambda p: transcribe_audio.invoke(p),
        '.zip': lambda p: extract_zip.invoke(p),
        '.jpg': lambda p: describe_image.invoke(p),
        '.jpeg': lambda p: describe_image.invoke(p),
        '.png': lambda p: describe_image.invoke(p),
    }
    
    processor = processors.get(ext)
    
    if processor:
        return processor(file_path)
    
    return f"[Unsupported file type: {ext}]"


# ============================================
# List of all file processing tools
# ============================================

file_processing_tools = [
    read_pdf,
    read_docx,
    read_pptx,
    read_text_file,
    read_csv,
    read_excel,
    read_jsonld,
    read_pdb,
    transcribe_audio,
    read_python_file,
    extract_zip,
    describe_image,
    read_file,
]
